from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


SRC = Path(r"C:\Users\Admin\Downloads\formslidebocobkn-250926141142-35737795.pptx")
OUT = Path(r"C:\Users\Admin\Downloads\SkinAI_slide_phan_3_ung_dung_AI.pptx")
IMG = Path(r"E:\DHBKDN\CapstoneProject\Webapp\SkinDeseases-FE\.tmp_report_images")

prs = Presentation(SRC)
W, H = prs.slide_width, prs.slide_height

TEAL = RGBColor(0, 128, 128)
DARK = RGBColor(17, 24, 39)
MUTED = RGBColor(82, 99, 118)
BORDER = RGBColor(214, 226, 226)
CARD = RGBColor(255, 255, 255)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(5, 150, 105)
AMBER = RGBColor(245, 158, 11)


def clear_slide(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 253, 253)
    bg.line.fill.background()

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.68))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(255, 255, 255)
    band.line.color.rgb = BORDER


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    bold=False,
    color=DARK,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.38), Inches(0.13), Inches(0.40), Inches(0.40)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background()

    add_text(
        slide,
        "3",
        Inches(0.38),
        Inches(0.18),
        Inches(0.40),
        Inches(0.22),
        12,
        True,
        RGBColor(255, 255, 255),
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "ỨNG DỤNG AI TRONG BÀI TOÁN PHÂN TÍCH ẢNH",
        Inches(0.88),
        Inches(0.16),
        Inches(7.1),
        Inches(0.28),
        13,
        True,
        TEAL,
    )
    add_text(slide, title, Inches(0.55), Inches(0.82), Inches(8.8), Inches(0.45), 24, True)
    if subtitle:
        add_text(
            slide,
            subtitle,
            Inches(0.58),
            Inches(1.24),
            Inches(8.8),
            Inches(0.34),
            10.5,
            False,
            MUTED,
        )


def add_footer(slide, page):
    add_text(
        slide,
        str(page),
        Inches(12.62),
        Inches(7.08),
        Inches(0.35),
        Inches(0.16),
        8,
        True,
        MUTED,
        PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        "SkinAI - Hệ thống hỗ trợ chẩn đoán bệnh da liễu từ hình ảnh",
        Inches(0.55),
        Inches(7.08),
        Inches(6.3),
        Inches(0.16),
        8,
        False,
        MUTED,
    )


def add_card(slide, x, y, w, h):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)
    return card


def add_bullets(slide, items, x, y, w, h, size=13, color=DARK):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(5)
    return box


def add_image_fit(slide, path, x, y, w, h):
    from PIL import Image

    im = Image.open(path)
    iw, ih = im.size
    scale = min(w / iw, h / ih)
    nw, nh = int(iw * scale), int(ih * scale)
    left = x + (w - nw) // 2
    top = y + (h - nh) // 2
    pic = slide.shapes.add_picture(str(path), left, top, width=nw, height=nh)

    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    rect.fill.background()
    rect.line.color.rgb = BORDER
    rect.line.width = Pt(1)
    return pic


def add_metric(slide, label, value, x, y, w, h, color=TEAL):
    add_card(slide, x, y, w, h)
    add_text(
        slide,
        value,
        x + Inches(0.08),
        y + Inches(0.10),
        w - Inches(0.16),
        Inches(0.28),
        18,
        True,
        color,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        label,
        x + Inches(0.08),
        y + Inches(0.45),
        w - Inches(0.16),
        Inches(0.22),
        8.5,
        False,
        MUTED,
        PP_ALIGN.CENTER,
    )


def add_section_chip(slide, text, x, y):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.35), Inches(0.34))
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(233, 252, 249)
    chip.line.color.rgb = RGBColor(157, 222, 215)
    add_text(slide, text, x + Inches(0.15), y + Inches(0.08), Inches(2.0), Inches(0.14), 8.8, True, TEAL)


def build_slide_12(slide):
    clear_slide(slide)
    add_bg(slide)
    add_title(
        slide,
        "Phát biểu bài toán AI",
        "Từ ảnh vùng da đầu vào, hệ thống trích xuất vùng tổn thương và phân loại nhóm bệnh phù hợp nhất.",
    )
    steps = [
        ("Ảnh đầu vào", "Người dùng tải ảnh vùng da nghi ngờ"),
        ("Kiểm tra ảnh", "Xác nhận ảnh có phải vùng da hợp lệ"),
        ("Segmentation", "Khoanh vùng tổn thương bằng U-Net"),
        ("Classification", "Phân loại 10 nhóm bệnh bằng CNN"),
        ("Kết quả & chatbot", "Hiển thị trực quan và tư vấn tham khảo"),
    ]
    x0, y, gap, cw, ch = Inches(0.55), Inches(2.0), Inches(0.17), Inches(2.35), Inches(1.35)
    for i, (title, desc) in enumerate(steps):
        x = x0 + i * (cw + gap)
        add_card(slide, x, y, cw, ch)
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.08), y + Inches(0.08), Inches(0.42), Inches(0.42)
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = TEAL
        circ.line.fill.background()
        add_text(slide, str(i + 1), x + Inches(0.08), y + Inches(0.16), Inches(0.42), Inches(0.18), 10, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.18), y + Inches(0.58), cw - Inches(0.36), Inches(0.25), 13, True, DARK, PP_ALIGN.CENTER)
        add_text(slide, desc, x + Inches(0.16), y + Inches(0.92), cw - Inches(0.32), Inches(0.30), 8.8, False, MUTED, PP_ALIGN.CENTER)
        if i < 4:
            add_text(slide, "→", x + cw + Inches(0.02), y + Inches(0.48), gap, Inches(0.20), 18, True, TEAL, PP_ALIGN.CENTER)
    add_section_chip(slide, "Vai trò của AI", Inches(0.75), Inches(4.08))
    add_bullets(
        slide,
        [
            "Tập trung vào vùng tổn thương thay vì toàn bộ ảnh, giúp giảm nhiễu nền.",
            "Kết hợp segmentation và classification để tạo kết quả trực quan, dễ giải thích.",
            "Kết quả AI được chuyển thành medical context cho chatbot tư vấn an toàn.",
        ],
        Inches(0.85),
        Inches(4.58),
        Inches(6.1),
        Inches(1.45),
        12.5,
    )
    add_card(slide, Inches(7.35), Inches(4.16), Inches(4.8), Inches(1.65))
    add_text(slide, "Đầu ra chính", Inches(7.6), Inches(4.38), Inches(4.2), Inches(0.25), 13, True, TEAL)
    add_bullets(
        slide,
        ["Mask phân đoạn vùng tổn thương", "ROI vùng bệnh sau segmentation", "Top-k bệnh da dự đoán và độ chính xác"],
        Inches(7.6),
        Inches(4.78),
        Inches(4.2),
        Inches(0.8),
        11.5,
    )
    add_footer(slide, 12)


def build_slide_13(slide):
    clear_slide(slide)
    add_bg(slide)
    add_title(
        slide,
        "Dữ liệu sử dụng cho mô hình AI",
        "Dữ liệu được kết hợp từ nhiều nguồn để phục vụ hai bài toán: phân đoạn tổn thương và phân loại bệnh da.",
    )
    add_card(slide, Inches(0.55), Inches(1.68), Inches(5.85), Inches(4.8))
    add_text(slide, "Segmentation - ISIC 2018 Task 1", Inches(0.8), Inches(1.9), Inches(5.3), Inches(0.24), 14, True, TEAL)
    add_image_fit(slide, IMG / "p56_img1_xref424.jpeg", Inches(0.82), Inches(2.28), Inches(5.3), Inches(2.55))
    add_bullets(
        slide,
        ["Ảnh dermoscopic đi kèm ground truth mask.", "Ảnh và mask được resize về 256×256.", "Dùng để huấn luyện mô hình khoanh vùng tổn thương."],
        Inches(0.9),
        Inches(4.98),
        Inches(5.0),
        Inches(1.1),
        10.8,
    )
    add_card(slide, Inches(6.75), Inches(1.68), Inches(5.65), Inches(4.8))
    add_text(slide, "Classification - 10 nhóm bệnh da", Inches(7.0), Inches(1.9), Inches(5.0), Inches(0.24), 14, True, TEAL)
    add_image_fit(slide, IMG / "p67_img1_xref447.png", Inches(7.05), Inches(2.28), Inches(5.0), Inches(2.55))
    add_bullets(
        slide,
        ["Tổng hợp từ HAM10000, ISIC, DermNet, PAD-UFES.", "Gồm 11.894 ảnh, chia train/validation/test.", "Bao phủ 10 lớp bệnh đã mapping trong hệ thống."],
        Inches(7.05),
        Inches(4.98),
        Inches(5.0),
        Inches(1.1),
        10.8,
    )
    add_footer(slide, 13)


def build_slide_14(slide):
    clear_slide(slide)
    add_bg(slide)
    add_title(slide, "Mô hình phân đoạn vùng tổn thương", "Sử dụng U-Net tùy chỉnh với encoder EfficientNet-B3 pretrained để tạo mask vùng bệnh.")
    add_card(slide, Inches(0.55), Inches(1.58), Inches(6.5), Inches(4.95))
    add_image_fit(slide, IMG / "p54_img1_xref418.png", Inches(0.82), Inches(1.85), Inches(5.95), Inches(3.65))
    add_text(
        slide,
        "Kiến trúc U-Net: encoder-decoder kết hợp skip connection để giữ thông tin biên và vị trí tổn thương.",
        Inches(0.9),
        Inches(5.68),
        Inches(5.8),
        Inches(0.35),
        10.5,
        False,
        MUTED,
        PP_ALIGN.CENTER,
    )
    add_card(slide, Inches(7.35), Inches(1.58), Inches(4.9), Inches(4.95))
    add_text(slide, "Cấu hình huấn luyện", Inches(7.65), Inches(1.88), Inches(4.3), Inches(0.25), 15, True, TEAL)
    config = [
        ("Model", "Custom U-Net + EfficientNet-B3"),
        ("Input", "256×256"),
        ("Output", "1 channel mask"),
        ("Batch size", "16"),
        ("Optimizer", "AdamW"),
        ("Loss", "0.7 Dice + 0.3 BCE"),
        ("Scheduler", "OneCycleLR"),
    ]
    y = Inches(2.35)
    for key, value in config:
        add_text(slide, key, Inches(7.75), y, Inches(1.4), Inches(0.18), 9.8, True, MUTED)
        add_text(slide, value, Inches(9.05), y, Inches(2.85), Inches(0.18), 9.8, False, DARK)
        y += Inches(0.38)
    add_text(
        slide,
        "Mục tiêu: xác định chính xác vùng tổn thương để hỗ trợ bước crop ROI và phân loại.",
        Inches(7.7),
        Inches(5.35),
        Inches(4.15),
        Inches(0.55),
        11.2,
    )
    add_footer(slide, 14)


def build_slide_15(slide):
    clear_slide(slide)
    add_bg(slide)
    add_title(slide, "Kết quả phân đoạn vùng tổn thương", "Mô hình đạt kết quả ổn định trên tập test, mask dự đoán bám sát vùng tổn thương thực tế.")
    add_card(slide, Inches(0.55), Inches(1.58), Inches(5.55), Inches(4.9))
    add_text(slide, "Đường cong huấn luyện", Inches(0.82), Inches(1.84), Inches(5.0), Inches(0.24), 13.5, True, TEAL)
    add_image_fit(slide, IMG / "p61_img1_xref437.png", Inches(0.82), Inches(2.2), Inches(5.0), Inches(3.42))
    add_text(slide, "Loss giảm ổn định; Dice validation duy trì quanh 0.90-0.91.", Inches(0.9), Inches(5.8), Inches(4.9), Inches(0.28), 10.2, False, MUTED, PP_ALIGN.CENTER)
    add_card(slide, Inches(6.42), Inches(1.58), Inches(5.95), Inches(4.9))
    add_text(slide, "Trực quan trên tập test", Inches(6.72), Inches(1.84), Inches(5.3), Inches(0.24), 13.5, True, TEAL)
    add_image_fit(slide, IMG / "p63_img2_xref441.jpeg", Inches(6.75), Inches(2.2), Inches(5.25), Inches(3.42))
    add_metric(slide, "Dice Score", "0.9094", Inches(0.75), Inches(6.55), Inches(1.55), Inches(0.78), GREEN)
    add_metric(slide, "IoU", "0.8376", Inches(2.48), Inches(6.55), Inches(1.55), Inches(0.78), TEAL)
    add_metric(slide, "Precision", "0.9259", Inches(4.21), Inches(6.55), Inches(1.55), Inches(0.78), BLUE)
    add_metric(slide, "Recall TTA", "0.9212", Inches(5.94), Inches(6.55), Inches(1.55), Inches(0.78), AMBER)
    add_text(slide, "TTA giúp tăng Recall, hỗ trợ phát hiện nhiều vùng tổn thương hơn trong suy luận.", Inches(7.8), Inches(6.72), Inches(4.3), Inches(0.35), 10.5, False, MUTED)
    add_footer(slide, 15)


def build_slide_16(slide):
    clear_slide(slide)
    add_bg(slide)
    add_title(slide, "Mô hình phân loại bệnh da", "EfficientNet-B0 được fine-tuning cho bài toán phân loại 10 nhóm bệnh da liễu.")
    add_card(slide, Inches(0.55), Inches(1.58), Inches(4.15), Inches(4.9))
    add_text(slide, "Dữ liệu và phân bố lớp", Inches(0.82), Inches(1.84), Inches(3.5), Inches(0.24), 13.5, True, TEAL)
    add_image_fit(slide, IMG / "p71_img1_xref456.png", Inches(0.85), Inches(2.25), Inches(3.55), Inches(2.1))
    add_bullets(
        slide,
        ["10 lớp bệnh da liễu.", "Có mất cân bằng lớp, NEVUS chiếm nhiều nhất.", "Áp dụng sampling và class-balanced loss."],
        Inches(0.92),
        Inches(4.65),
        Inches(3.45),
        Inches(1.25),
        10,
    )
    add_card(slide, Inches(4.95), Inches(1.58), Inches(3.95), Inches(4.9))
    add_text(slide, "Kiến trúc EfficientNet-B0", Inches(5.2), Inches(1.84), Inches(3.45), Inches(0.24), 13.5, True, TEAL)
    add_image_fit(slide, IMG / "p64_img1_xref443.jpeg", Inches(5.38), Inches(2.22), Inches(3.18), Inches(2.65))
    add_bullets(
        slide,
        ["Pretrained ImageNet.", "Thay classifier cho 10 lớp.", "Cân bằng giữa độ chính xác và chi phí tính toán."],
        Inches(5.18),
        Inches(5.04),
        Inches(3.45),
        Inches(0.9),
        10,
    )
    add_card(slide, Inches(9.15), Inches(1.58), Inches(3.25), Inches(4.9))
    add_text(slide, "Cấu hình chính", Inches(9.42), Inches(1.84), Inches(2.6), Inches(0.24), 13.5, True, TEAL)
    config = [
        ("Input", "256×256"),
        ("Batch", "64"),
        ("Epoch", "100"),
        ("Optimizer", "AdamW"),
        ("Scheduler", "OneCycleLR"),
        ("Weight decay", "5e-4"),
        ("TTA", "Có"),
    ]
    y = Inches(2.35)
    for key, value in config:
        add_text(slide, key, Inches(9.45), y, Inches(1.05), Inches(0.16), 9.5, True, MUTED)
        add_text(slide, value, Inches(10.55), y, Inches(1.5), Inches(0.16), 9.5)
        y += Inches(0.35)
    add_text(slide, "Pipeline giữ nguyên preprocessing, augmentation, sampler, loss và evaluation để đánh giá công bằng.", Inches(9.42), Inches(5.25), Inches(2.55), Inches(0.6), 9.5, False, MUTED)
    add_footer(slide, 16)


def build_slide_17(slide):
    clear_slide(slide)
    add_bg(slide)
    add_title(slide, "Kết quả phân loại và so sánh mô hình", "Mô hình phân loại được đánh giá bằng Accuracy, Macro-F1, Weighted-F1 và ma trận nhầm lẫn trên tập test.")
    add_card(slide, Inches(0.55), Inches(1.55), Inches(3.7), Inches(4.75))
    add_text(slide, "Classification report", Inches(0.8), Inches(1.8), Inches(3.2), Inches(0.24), 13, True, TEAL)
    add_image_fit(slide, IMG / "p75_img2_xref470.png", Inches(0.86), Inches(2.18), Inches(3.1), Inches(2.25))
    add_metric(slide, "Test Accuracy", "75%", Inches(0.85), Inches(4.85), Inches(1.45), Inches(0.70), TEAL)
    add_metric(slide, "Macro-F1", "0.74", Inches(2.48), Inches(4.85), Inches(1.45), Inches(0.70), GREEN)
    add_text(slide, "ACNE, BCC và NEVUS đạt F1-score tốt nhất trong tập test.", Inches(0.9), Inches(5.75), Inches(3.1), Inches(0.32), 9.5, False, MUTED, PP_ALIGN.CENTER)

    add_card(slide, Inches(4.55), Inches(1.55), Inches(3.55), Inches(4.75))
    add_text(slide, "Ma trận nhầm lẫn", Inches(4.82), Inches(1.8), Inches(3.0), Inches(0.24), 13, True, TEAL)
    add_image_fit(slide, IMG / "p76_img1_xref472.png", Inches(4.88), Inches(2.18), Inches(2.95), Inches(3.05))
    add_text(slide, "Giúp quan sát các cặp lớp dễ nhầm lẫn khi phân loại.", Inches(4.9), Inches(5.62), Inches(2.85), Inches(0.28), 9.5, False, MUTED, PP_ALIGN.CENTER)

    add_card(slide, Inches(8.42), Inches(1.55), Inches(3.95), Inches(4.75))
    add_text(slide, "So sánh EfficientNet-B0 và ResNet-50", Inches(8.68), Inches(1.8), Inches(3.4), Inches(0.24), 12.5, True, TEAL)
    rows = [
        ("Tham số", "4.67M", "24.56M"),
        ("Train time", "69.4 phút", "110.7 phút"),
        ("Test Acc TTA", "74.89%", "77.85%"),
        ("Test Macro-F1 TTA", "73.77%", "77.43%"),
    ]
    add_text(slide, "Chỉ số", Inches(8.65), Inches(2.35), Inches(1.2), Inches(0.18), 9.2, True, MUTED)
    add_text(slide, "Eff-B0", Inches(9.9), Inches(2.35), Inches(0.9), Inches(0.18), 9.2, True, MUTED, PP_ALIGN.CENTER)
    add_text(slide, "ResNet", Inches(11.0), Inches(2.35), Inches(0.9), Inches(0.18), 9.2, True, MUTED, PP_ALIGN.CENTER)
    y = Inches(2.78)
    for key, eff, res in rows:
        add_text(slide, key, Inches(8.65), y, Inches(1.2), Inches(0.18), 8.9)
        add_text(slide, eff, Inches(9.85), y, Inches(0.95), Inches(0.18), 8.9, True, TEAL, PP_ALIGN.CENTER)
        add_text(slide, res, Inches(10.95), y, Inches(0.95), Inches(0.18), 8.9, True, BLUE, PP_ALIGN.CENTER)
        y += Inches(0.45)
    add_text(slide, "Kết luận", Inches(8.65), Inches(4.78), Inches(3.2), Inches(0.22), 11, True)
    add_bullets(
        slide,
        ["ResNet-50 tốt hơn về độ chính xác cuối cùng.", "EfficientNet-B0 gọn nhẹ hơn, phù hợp triển khai VPS CPU."],
        Inches(8.72),
        Inches(5.14),
        Inches(3.25),
        Inches(0.75),
        9.5,
    )
    add_footer(slide, 17)


builders = [build_slide_12, build_slide_13, build_slide_14, build_slide_15, build_slide_16, build_slide_17]
for offset, builder in enumerate(builders, start=11):
    builder(prs.slides[offset])

prs.save(OUT)
print(f"Saved: {OUT}")
